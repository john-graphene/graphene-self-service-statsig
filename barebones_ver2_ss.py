# import warnings
# warnings.filterwarnings('ignore')  
import os
import io
import shutil
import pandas as pd
import numpy as np
import regex as re
import json
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
# from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR, PP_ALIGN
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from docx import Document
"""""""""""""""""""""""""""""""""""""""  Inputs """""""""""""""""""""""""""""""""""""""""""""

## REVISED GLOBAL VARS, PENDING FRONTEND INTEGRATION
default_template = "./Files/PPT_Template/Proposed_Barebones.pptx"
placeholder_description_input = "./Files/PPT_Template/Ref material/placeholder_description_input.csv"
slide_master_placeholders_output = "./Files/PPT_Template/Ref material/Slide_Master_Placeholders_output.csv"
table_config_json = './Files/PPT_Template/Ref material/table_format_json.txt'
blank_mapping_path = "./Files/PPT_Template/mapping_blank.xlsx"
encoding = 'utf-8-sig'

project_name = "Friska"
# project_name = "Salonpas_Samples"

brief_doc_path = f"./Files/{project_name}/AI_brief.docx"
mapping_path = f"./Files/{project_name}/mapping_input.xlsx"
workbook_path = f"./Files/{project_name}/workbook.xlsx"
drivers_equity_path = f"./Files/{project_name}/drivers.csv"
output_file_name = f"./Files/{project_name}/barebones.pptx"

"""""""""""""""""""""""""""""""""""""""  Generate Mapping File """""""""""""""""""""""""""""""""""""""""""""


def gen_brief_df(brief_doc_path):
    """
        Generate combined table based on AI brief, for use in mapping df //
        WORK IN PROGRESS: PENDING AI BRIEF TEMPLATE CONFIRMATION
    """
    docx = Document(brief_doc_path)
    brief_df = pd.DataFrame()

    ## Get all tables in AI Brief which contains questions
    for table in docx.tables:
        first_line_merge_condition = (
                    (table.cell(0, 0).text == table.cell(0, 1).text) | (table.cell(0, 1).text == table.cell(0, 2).text))
        if (table.cell(0, 0).text == "Questions to answer") | (
                first_line_merge_condition & (table.cell(1, 0).text == "Questions to answer")):
            for row in table.rows:
                temp_list = []
                for cell in row.cells:
                    temp_list.append(cell.text)
                temp_df = pd.DataFrame(temp_list).T
                brief_df = pd.concat((brief_df, temp_df), axis=0)

    ## Concat all tables into 1 table, rename column header
    brief_df = brief_df.drop_duplicates().reset_index(drop=True)
    indices_to_drop = []
    for idx, row in brief_df.iterrows():
        if row[0] == row[1] == row[2]:
            indices_to_drop.append(idx)
    brief_df = brief_df.drop(indices_to_drop)
    brief_df.columns = brief_df.iloc[0].values
    brief_df = brief_df.iloc[1:].reset_index(drop=True)
    brief_df.columns = [col.split(" ")[0] for col in brief_df.columns]

    def explode_series(brief_df, col, split_substring="\n\n"):
        """ Splits string column by specified substring across multiple rows, replicates other columns """
        brief_df = brief_df.reset_index(drop=True)  ## necessary for pd.explode to function correctly
        empty_list = []
        for text in brief_df[col]:
            text_list = text.strip().split(split_substring)
            empty_list.append(text_list)
        brief_df[col] = pd.Series(empty_list)
        brief_df = brief_df.explode(col, ignore_index=True)
        brief_df = brief_df[brief_df[col] != '']
        return brief_df

    ## Split Insights column with multiple asks
    brief_df = explode_series(brief_df, "Insights", "\n")
    brief_df['Insights'] = brief_df['Insights'].str.replace('%', 'Size')
    brief_df = brief_df[~(brief_df['Insights'].str.contains('AI Discovered|Predefined|Suggestions', case=False))]

    ## Retain Questions column with '?'
    brief_df = explode_series(brief_df, "Questions", "\n\n")
    brief_df = brief_df[brief_df['Questions'].str.contains(r'\?')]

    return brief_df.reset_index(drop=True)


def gen_mapping_file(workbook_path, brief_doc_path, drivers_equity_path, blank_mapping_path, mapping_path):
    """
        Generate default mapping file based input XLS file for download, for user ref //
        Also fills in Questions to answer from AI brief //
        Return mapping df //
        WORK IN PROGRESS: PENDING MAPPING FILE CONFIRMATION INCL. STATSIG AND FRONT END INTEGRATION
    """
    xls_dict = pd.read_excel(workbook_path, sheet_name=None)
    brief_df = gen_brief_df(brief_doc_path)

    intro_slides_df = pd.DataFrame([["Deck_Cover"], ["Project_Objective"], ["Project_Scope"]], columns=["Slide_Type"])

    base_df = xls_dict['Base']
    if len(base_df['Brand'].unique()) > 1:
        base_df = base_df[(~base_df['No_of_Conversations'].isna()) & (base_df['Brand'] == "All")].copy()
    else:
        base_df = base_df[(~base_df['No_of_Conversations'].isna())].copy()
    base_df['Slide_Type'] = "SOW"
    base_df = base_df[['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Slide_Type']].drop_duplicates()

    performance_df = xls_dict['Performance']
    performance_df['Slide_Type'] = "Performance"
    performance_df = performance_df[
        ['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type', 'Subtype', 'Slide_Type']].drop_duplicates()

    de_df = None
    try:
        de_df = pd.read_csv(drivers_equity_path)
        de_df = de_df.rename(columns={'SubCategory': 'Subcategory'}, errors='ignore')
        #         de_df = de_df[de_df['Brand']=="All"].copy()
        if len(de_df['Brand'].unique()) > 1:
            de_df['Brand'] = "All"
        de_df['Type'] = "Drivers & Equity"
        de_df['Subtype'] = "Of purchase/usage"
        de_df['Slide_Type'] = "Drivers"
        de_df = de_df[['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type', 'Subtype',
                       'Slide_Type']].drop_duplicates()
    except:
        pass

    kol_df = None
    if 'KOLs Influencers' in xls_dict.keys():
        kol_df = xls_dict['KOLs Influencers']
        kol_df['Slide_Type'] = "KOL"
        kol_df = kol_df[['Segment', 'Slide_Type']].drop_duplicates()

    et_df = None
    if 'Emerging Trends' in xls_dict.keys():
        et_df = xls_dict['Emerging Trends']
        et_df.columns = [col.replace("Key", "") for col in et_df.columns]
        et_df = et_df.rename(columns={'SubCategory': 'Subcategory'}, errors='ignore')
        et_df['Subtype'] = "Across all"  ## PENDING CONFIRMATION
        et_df['Slide_Type'] = "Emerging Trends"
        et_df = et_df[['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type', 'Subtype',
                       'Slide_Type']].drop_duplicates()

    mapping_df = pd.read_excel(blank_mapping_path, sheet_name="Mapping")
    mapping_df = pd.concat((mapping_df, intro_slides_df), axis=0)
    mapping_df = pd.concat((mapping_df, base_df), axis=0)
    mapping_df = pd.concat((mapping_df, performance_df), axis=0)
    try:
        mapping_df = pd.concat((mapping_df, de_df), axis=0)
    except:
        pass
    try:
        mapping_df = pd.concat((mapping_df, et_df), axis=0)
    except:
        pass
    try:
        mapping_df = pd.concat((mapping_df, kol_df), axis=0)
    except:
        pass
    mapping_df = mapping_df.reset_index(drop=True)

    for idx, row in mapping_df.iterrows():
        if pd.isna(row['Type']) == False:
            performance_type = row['Type']
            performance_subtype = row['Subtype']
            sub_brief_df = brief_df.loc[(brief_df['Insights'].str.contains(performance_type, case=False))]
            #             print(f"Type:{performance_type}\nSubtype:{performance_subtype}\nLen:{len(sub_brief_df)}")
            try:
                if len(sub_brief_df) == 1:
                    row['Questions'] = sub_brief_df['Questions'].item()
                else:
                    try:
                        sub_sub_brief_df = sub_brief_df.loc[
                            (sub_brief_df['Insights'].str.contains(performance_subtype, case=False))]
                        row['Questions'] = sub_sub_brief_df['Questions'].head(1).item()
                    except ValueError:
                        row['Questions'] = sub_brief_df['Questions'].head(1).item()
            except ValueError:
                pass

    output_path = shutil.copy(blank_mapping_path, mapping_path)
    with pd.ExcelWriter(output_path,
                        engine='openpyxl',
                        mode='a',
                        if_sheet_exists='overlay') as writer:
        mapping_df.to_excel(writer, sheet_name="Mapping", index=False)

    return mapping_df


def get_mapping_df(workbook_path, brief_doc_path, drivers_equity_path, blank_mapping_path, mapping_path):
    """
        Get mapping table based on uploaded mapping file with user (BS) inputs //
        WORK IN PROGRESS: PENDING FRONTEND INTEGRATION
    """
    try:
        mapping_df = pd.read_excel(mapping_path, sheet_name="Mapping").copy()
        print("Mapping file retrieved.")
        return mapping_df
    except FileNotFoundError:
        print("No mapping file input from user, generating default mapping file.")
        default_mapping_df = gen_mapping_file(workbook_path, brief_doc_path, drivers_equity_path, blank_mapping_path, mapping_path)
        return default_mapping_df


def gen_sequence_dict(mapping_df):
    """
        Gen slide sequence dict based on mapping df //
        WORK IN PROGRESS: INTEGRATE GENERATION OF COVER, PROJECT_OBJECTIVE, PROJECT_SCOPE SLIDES
    """
    ## slide_sequence_dict keys
    slide_sequence_keys = []
    slide_sequence_list = list(mapping_df['Slide_Type'])
    type_cnt = 1
    for idx, slide_type in enumerate(slide_sequence_list):
        if slide_type == slide_sequence_list[idx - 1]:
            type_cnt += 1
        else:
            type_cnt = 1
        slide_sequence_keys.append(f"{slide_type} {type_cnt}")
    #     print(slide_sequence_keys)
    mapping_df['Slide_Type'] = slide_sequence_keys

    ## slide_sequence_dict values
    slide_sequence_dict = {}
    for idx, row in mapping_df.iterrows():
        slide_type = row['Slide_Type']
        row = row.drop('Slide_Type')
        row = row.dropna()
        slide_sequence_dict[slide_type] = pd.DataFrame(row).T

    return slide_sequence_dict


"""""""""""""""""""""""""""""""""""""""  Excel Data Processing """""""""""""""""""""""""""""""""""""""""""""



def get_sow_data(base_df, master_df, slide_sequence_dict):
    """
        Get data from Base and Master sheets for Scope of Work slide
    """
    if len(base_df['Brand'].unique()) > 1:
        sow_df = base_df[(~base_df['No_of_Conversations'].isna())&(base_df['Brand']=="All")].copy()
    else:
        sow_df = base_df[(~base_df['No_of_Conversations'].isna())].copy()
#     sow_df = base_df[(~base_df['No_of_Conversations'].isna())&(base_df['Brand']=="All")]
    sow_df = sow_df.merge(master_df, on='PeriodKey').reset_index(drop=True)
    sow_df = sow_df[['Category','Subcategory','Country','Segment','Brand','No_of_People','No_of_Conversations','Cycle']]
    sow_df['Slide_Type'] = "SOW"
    for key, value in slide_sequence_dict.items():
        if "SOW" in key:
            merged_df = value.merge(sow_df, on = ['Category','Subcategory','Country','Segment','Brand'], how = 'left')
            slide_sequence_dict[key] = merged_df
    return slide_sequence_dict


def get_performance_data(performance_df, base_df, slide_sequence_dict):
    """
        Get data from Performance and Base sheets for Performance slide //
        WORK IN PROGRESS: CROSS SEGMENT, STATSIG COMPUTATION
    """
    base_df_trimmed = base_df.drop(
        columns=['PeriodDateEnd', 'Month', 'Quarter', 'Year', 'Trial_Percent', 'No_of_Conversations',
                 'Comments']).drop_duplicates().reset_index(drop=True)
    performance_df_trimmed = performance_df.drop(columns=['Denominator', 'Comments']).drop_duplicates().reset_index(
        drop=True)
    pb_df = performance_df_trimmed.merge(base_df_trimmed,
                                         on=['PeriodKey', 'Category', 'Subcategory', 'Country', 'Segment', 'Brand'],
                                         how='left')
    pb_df['No_of_People'] = pb_df['No_of_People'].fillna(-1)

    pb_df = pb_df.convert_dtypes()
    pb_df['Measure Value'] = pb_df['Measure Value'].astype(object)
    pb_df['Content'] = pb_df['Content'].apply(lambda x: str(x)[0].capitalize() + str(x)[1:])
    pb_df.loc[pb_df['Measure Type'] == "%", "Measure Value"] = pb_df['Measure Value'].apply(lambda x: f"{x / 100:.0%}")
    pb_df = pb_df[
        ['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type', 'Subtype', 'Content', 'Measure Value',
         'No_of_People']]

    segment_separator = "//"
    for key, slide_df in slide_sequence_dict.items():
        if "Performance" in key:
            ## Single Segment
            if not slide_df['Segment'].str.contains(segment_separator).any():
                merged_df = slide_df.merge(pb_df, on=['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type',
                                                      'Subtype'], how='left')
                merged_df['Slide_Type'] = "Performance"
                merged_df = merged_df.sort_values('Measure Value', ascending=False).reset_index(drop=True)
                slide_sequence_dict[key] = merged_df
            ## Cross Segment
            if slide_df['Segment'].str.contains(segment_separator).any():
                ## split Series 'Segment' to DataFrame with columns 'Segment 1', 'Segment 2', etc.
                split_segment = slide_df['Segment'].str.split(segment_separator, expand=True)
                split_segment.columns = [col_name if type(col_name) == "str" else f"Segment {col_name + 1}" for col_name
                                         in split_segment.columns]

                ## filter performance data on slide-specific cross segments to avoid duplication with similar Type/Subtype
                temp_pb_df = pb_df[pb_df['Segment'].isin(split_segment.values[0])]

                ## drop 'Segment' column + potentially unfilled columns
                col_drop_list = ['Segment', 'Questions', 'Statsig_Type', 'Ref_Benchmark']
                temp_slide_df = slide_df.drop(columns=col_drop_list, errors='ignore')
                temp_slide_df = temp_slide_df.merge(temp_pb_df, how='left',
                                                    on=['Category', 'Subcategory', 'Country', 'Brand', 'Type',
                                                        'Subtype'])
                temp_slide_df = temp_slide_df.groupby(
                    ['Category', 'Subcategory', 'Country', 'Brand', 'Type', 'Subtype', 'Content',
                     'Segment']).sum().unstack().reset_index()

                ## rename cross-segment values in column names to 'Segment 1/2/etc'
                split_segment_dict = {}
                for col in split_segment.columns:
                    split_segment_dict[split_segment[col].values[0]] = col
                    temp_slide_df[col] = split_segment[col].values[0]
                temp_slide_df = temp_slide_df.rename(columns=split_segment_dict)
                temp_slide_df.columns = [' '.join(col).strip() for col in temp_slide_df.columns.values]
                temp_slide_df = temp_slide_df.sort_values(['Measure Value Segment 1', 'Measure Value Segment 2'],
                                                          ascending=False, na_position='last').reset_index(drop=True)

                ## add back potentially unfilled columns
                for col in col_drop_list[1:]:
                    try:
                        temp_slide_df[col] = slide_df[col].unique()[0]
                    except:
                        pass
                temp_slide_df['Slide_Type'] = "Performance"
                slide_sequence_dict[key] = temp_slide_df
    return slide_sequence_dict


def get_drivers_data(drivers_equity_path, base_df, mapping_df, slide_sequence_dict):
    """
        Get data from Drivers & Equity CSV for Drivers slide //
        WORK IN PROGRESS: PENDING FRONTEND INTEGRATION
    """
    de_df = pd.read_csv(drivers_equity_path).copy()
    de_df = de_df.rename(columns={'SubCategory': 'Subcategory'}, errors="ignore")
    de_df = de_df.drop(columns=['PeriodDateEnd', 'Month', 'Quarter', 'Year', 'CBI'], errors="ignore")
    base_df_trimmed = base_df.drop(
        columns=['PeriodDateEnd', 'Month', 'Quarter', 'Year', 'Trial_Percent', 'No_of_Conversations', 'Comments',
                 'Validations'], errors="ignore").drop_duplicates().reset_index(drop=True)

    if len(base_df['Brand'].unique()) > 1:
        base_df_trimmed = base_df_trimmed[base_df_trimmed['Brand'] == "All"].copy()
    base_df_trimmed = base_df_trimmed.drop(columns=['Brand']).drop_duplicates()

    de_df = de_df.merge(base_df_trimmed, how="left", on=['PeriodKey', 'Category', 'Subcategory', 'Country', 'Segment'])

    de_df = de_df.drop(columns=['PeriodKey'], errors='ignore')

    for idx, segment in enumerate(mapping_df[mapping_df['Slide_Type'].str.contains("Drivers")]['Segment'].unique()):
        sub_de_df = de_df[de_df['Segment'] == segment].copy()
        sub_de_df = sub_de_df[sub_de_df['Brand'] != "All"]
        sub_de_df['Brand'] = sub_de_df['Brand'].apply(lambda x: str(x)[0].capitalize() + str(x)[1:])
        sub_de_df['Driver'] = sub_de_df['Driver'].apply(lambda x: str(x)[0].capitalize() + str(x)[1:])
        sub_de_df = sub_de_df.groupby(
            ['Category', 'Subcategory', 'Country', 'Segment', 'No_of_People', 'Driver_Score', 'Driver',
             'Brand']).sum().unstack().sort_values('Driver_Score', ascending=False).reset_index()
        sub_de_df.columns = [' '.join(col).strip().replace("Equity_Score ", "") for col in sub_de_df.columns.values]
        sub_de_df['Slide_Type'] = "Drivers"
        sub_de_df = sub_de_df.merge(slide_sequence_dict[f'Drivers {idx + 1}'],
                                    on=['Category', 'Subcategory', 'Country', 'Segment'],
                                    how='outer').drop(columns=['Brand'])
        slide_sequence_dict[f'Drivers {idx + 1}'] = sub_de_df
    return slide_sequence_dict


def get_et_data(base_df, emergingtrends_df, slide_sequence_dict):
    """
        Get data from Emerging Trends sheet for Emerging Trends slide //
        WORK IN PROGRESS: PENDING LOGIC CONFIRMATION, FIX UNPIVOT
    """
    base_df_trimmed = base_df.drop(
        columns=['PeriodDateEnd', 'Month', 'Quarter', 'Year', 'Trial_Percent', 'No_of_Conversations',
                 'Comments']).drop_duplicates().reset_index(drop=True)
    et_df_trimmed = emergingtrends_df.drop(columns=['PeriodDateEnd', 'Month']).drop_duplicates().reset_index(drop=True)
    et_df_trimmed = et_df_trimmed.rename(
        columns={'CategoryKey': 'Category', 'SubCategoryKey': 'Subcategory', 'CountryKey': 'Country',
                 'SegmentKey': 'Segment', 'BrandKey': 'Brand'}, errors="ignore")
    et_df = et_df_trimmed.merge(base_df_trimmed,
                                on=['PeriodKey', 'Category', 'Subcategory', 'Country', 'Segment', 'Brand'],
                                how='left')
    #     print(et_df.columns, et_df.shape)
    et_df = et_df.rename(columns={'Number of People Discussing': '#Discussing',
                                  'Positive People %': '%PosPpl',
                                  'TS Positive': 'TS_Pos',
                                  'No_of_People': '#Ppl'})
    et_df['Period'] = et_df['Quarter'].astype(str) + " " + et_df['Year'].astype(str)
    et_df = et_df.drop(columns=['Quarter', 'Year', '%PosPpl', 'TS_Pos',
                                '#Ppl'])  ## PENDING CONFIRMATION WHICH METRIC TO CHART, TO REVIEW
    #     print(et_df.columns)

    periodkey_list = et_df['PeriodKey'].unique()
    first_et_df = et_df[et_df['PeriodKey'] == min(periodkey_list)]
    last_et_df = et_df[et_df['PeriodKey'] == max(periodkey_list)]
    ends_et_df = first_et_df.merge(last_et_df,
                                   on=['Category', 'Subcategory', 'Country', 'Segment', 'Brand', 'Type', 'Subtype',
                                       'Content'],
                                   how="inner")
    #     print(ends_et_df.columns)
    ends_et_df['#Discussing_Var'] = ends_et_df['#Discussing_y'] - ends_et_df['#Discussing_x']
    pos_ends_et_df = ends_et_df[ends_et_df['#Discussing_Var'] > 0]
    pos_ends_et_df = pos_ends_et_df[
        ['Period_x', 'Period_y', 'Category', 'Subcategory', 'Country', 'Segment', 'Content', '#Discussing_x',
         '#Discussing_y']]
    pos_ends_et_df = pd.melt(pos_ends_et_df, id_vars=['Category', 'Subcategory', 'Country', 'Segment', 'Content'])

    pos_ends_et_df['Slide_Type'] = "Emerging Trends"  ##TBC
    slide_sequence_dict['Emerging Trends 1'] = pos_ends_et_df
    return slide_sequence_dict


def get_kol_data(kols_influencers_df, slide_sequence_dict):
    """
        Get data from KOL Influencers sheet for KOL slide //
        WORK IN PROGRESS: PENDING EMPTIES STRUCTURES CONFIRMATION
    """
    kol_df = kols_influencers_df.copy()
    ## drop empty columns
    drop_col = []
    for col in kol_df.columns:
        if len(kol_df[kol_df[col].isna()]) / len(kol_df[col]) > 0.9:
            drop_col.append(col)
    kol_df = kol_df.drop(columns=drop_col, errors="ignore")

    ## check if Impact Score in columns
    col_found = False
    idx_shift = 0
    while col_found == False:
        for col in kol_df.columns:
            if re.search("impact score", col, re.IGNORECASE):
                col_found = True
                break
        if col_found == False:
            kol_df.columns = kol_df.iloc[idx_shift]
            kol_df = kol_df.iloc[idx_shift + 1:]

    ## regex column titles to replace whitespace and remove parentheses
    new_cols = []
    for col in kol_df.columns:
        if (type(col) == str) & ("segment" in col.lower()):
            col = "Segment"
            new_cols.append(col)
        if (type(col) == str) & ("name" in col.lower()):
            col = "Name"
            new_cols.append(col)
        if (type(col) == str) & ("impact" in col.lower()):
            col = "Impact Score"
            new_cols.append(col)
        if (type(col) == str) & ("platform" in col.lower()):
            col = "Engagement Platform"
            new_cols.append(col)
        if (type(col) == str) & ("institution" in col.lower()):
            col = "Institution"
            new_cols.append(col)
    kol_df.columns = new_cols

    ## split to top and bottom tables and append to dict
    for slide_count, segment in enumerate(kol_df['Segment'].unique()):
        kol_sub_df = kol_df[kol_df['Segment'] == segment].reset_index(drop=True)
        kol_sub_df['Slide_Type'] = "KOL"
        if segment == slide_sequence_dict[f'KOL {slide_count + 1}']['Segment'].unique().item():
            slide_sequence_dict[f'KOL {slide_count + 1}'] = kol_sub_df

    return slide_sequence_dict


def get_slide_sequence_values(mapping_path):
    """
        Get all data as sequential dict //
        WORK IN PROGRESS: PENDING CHANGES TO PERFORMANCE/DRIVERS FUNCTIONS FOR STATSIG INTEGRATION
    """
    mapped_df = get_mapping_df(workbook_path, brief_doc_path, drivers_equity_path, blank_mapping_path, mapping_path)
    slide_sequence_dict = gen_sequence_dict(mapped_df)

    df_list = []
    ## to be adjusted to take input from uploaded file
    xls_dict = pd.read_excel(workbook_path, sheet_name=None)
    for sheet_name, workbook_df in xls_dict.items():
        df_name = str(re.sub(r'[\s+ %]', '', sheet_name.lower().strip())) + "_df"
        globals()[df_name] = workbook_df
        df_list.append(df_name)
    for df in df_list:
        exec(f"{df} = {df}.loc[:, ~{df}.columns.str.contains('^Unnamed')]")
        exec(f"{df} = {df}.convert_dtypes()")

    #     slide_layout_name_list = ['Cover','Objective','Scope']
    slide_type_list = ['Deck_Cover', 'Project_Objective', 'Project_Scope']
    for key, df in slide_sequence_dict.items():
        key_name = str(re.sub(r'[\d]', '', key)).strip()
        if key_name in slide_type_list:
            intro_df = pd.DataFrame([key_name], columns=["Slide_Type"])
            slide_sequence_dict[key] = intro_df

    slide_sequence_dict = get_sow_data(base_df, master_df, slide_sequence_dict)
    slide_sequence_dict = get_performance_data(performance_df, base_df, slide_sequence_dict)
    try:
        slide_sequence_dict = get_drivers_data(drivers_equity_path, base_df, mapped_df, slide_sequence_dict)
    except Exception as Err:
        print(Err)
        pass
    try:
        slide_sequence_dict = get_et_data(base_df, emergingtrends_df, slide_sequence_dict)
    except Exception as Err:
        print(Err)
        pass
    try:
        slide_sequence_dict = get_kol_data(kols_influencers_df, slide_sequence_dict)
    except Exception as Err:
        print(Err)
        pass
    return slide_sequence_dict


"""""""""""""""""""""""""""""""""""""""  PPT Template """""""""""""""""""""""""""""""""""""""""""""


def get_slide_master_df(default_template, placeholder_description_input, slide_master_placeholders_output):
    """
        Get Slide Layout Names based on PPTX Template //
        Template last updated 2023-09-26 //
        WORK IN PROGRESS: UPDATE TO TEMPLATE FOR DRIVERS SLIDE TEMPLATE & PENDING FRONTEND INTEGRATION
    """
    prs = Presentation(default_template)
    slide_master_df = pd.DataFrame(
        columns=["Slide_Master_Index", "Slide_Layout_Index", "Slide_Layout_Name", "Placeholder_Index",
                 "Placeholder_Name"])
    for sm_idx, sm in enumerate(prs.slide_masters):
        for sl_idx, sl in enumerate(sm.slide_layouts):
            for sh_idx, shape in enumerate(sl.shapes):
                if shape.is_placeholder:
                    slide_master_dict = {"Slide_Master_Index": sm_idx,
                                         "Slide_Layout_Index": sl_idx,
                                         "Slide_Layout_Name": sl.name,
                                         "Placeholder_Index": shape.placeholder_format.idx,
                                         "Placeholder_Name": shape.name
                                         }
                    slide_master_df = pd.concat((slide_master_df, pd.DataFrame([slide_master_dict])), axis=0)

    ## Merge Placeholder description
    placeholder_description_df = pd.read_csv(placeholder_description_input)
    slide_master_df = slide_master_df.merge(placeholder_description_df,
                                            on=['Slide_Layout_Index', 'Slide_Layout_Name', 'Placeholder_Index'],
                                            how='left')

    ## Merge Slide_Type
    slide_description_df = pd.DataFrame([["Cover", "Deck_Cover"],
                                         ["Objective", "Project_Objective"],
                                         ["Scope", "Project_Scope"],
                                         ["Divider_normal", "Section_Divider"],
                                         ["Scope_of_Work", "SOW"],
                                         ["Content_Vertical", "Performance"],
                                         ["Content_Horizontal_Chart", "Emerging Trends"],
                                         ["Content_Vertical", "Drivers"],
                                         ["KOL", "KOL"],
                                         ["End", "Deck_End"],
                                         ["blank", "blank_slide"]
                                         ], columns=['Slide_Layout_Name', 'Slide_Type'])
    slide_master_df = slide_master_df.merge(slide_description_df, on="Slide_Layout_Name", how="left").sort_values(
        ['Slide_Master_Index', 'Slide_Layout_Index', 'Placeholder_Index']).reset_index(drop=True)
    if not os.path.isfile(slide_master_placeholders_output):
        slide_master_df.to_csv(slide_master_placeholders_output, index=False)

    return slide_master_df


def fill_text_data(slide, slide_df, slide_master_df, brief_doc_path):
    """
        Fill Text Placeholders with required text inputs //
        WORK IN PROGRESS: PENDING CONFIRMATION OF AI BRIEF WITH BS, UPDATE FOR CROSS-SEGMENT
    """
    slide_type = slide_df['Slide_Type'].unique()[0]

    ## get PH index based on template, last updated 2023-09-26
    def get_ph_idx(slide_type, ph_desc):
        sub_slide_master_df = slide_master_df[slide_master_df['Slide_Type'] == slide_type]
        try:
            return sub_slide_master_df[sub_slide_master_df['Placeholder_Description'] == ph_desc][
                'Placeholder_Index'].unique()[0]
        except:
            return None

    ## fill Text PH with text content
    for shape in slide.shapes:
        if slide.slide_layout.name == "Cover":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Cover Text PH"):
                doc_header = Document(brief_doc_path).paragraphs[0].text
                shape.text = doc_header.title()
                para_2 = shape.text_frame.add_paragraph()
                run_2_0 = para_2.add_run()
                para_3 = shape.text_frame.add_paragraph()
                run_3_0 = para_3.add_run()
                para_4 = shape.text_frame.add_paragraph()
                run_4_0 = para_4.add_run()
                para_5 = shape.text_frame.add_paragraph()
                run_5_0 = para_5.add_run()
                run_5_1 = para_5.add_run()
                run_2_0.text = "<Report Title>"
                run_2_0.font.size = Pt(24)
                run_2_0.font.bold = False
                run_3_0.text = "<Mmmm YYYY>"
                run_3_0.font.size = Pt(16)
                run_3_0.font.bold = False
                run_4_0.text = " "
                run_4_0.font.size = Pt(36)
                run_4_0.font.bold = False
                run_5_0.text = "AI"
                run_5_1.text = " based research"
                run_5_0.font.size = Pt(20)
                run_5_1.font.size = Pt(20)
                run_5_0.font.color.rgb = RGBColor(255, 0, 0)
                run_5_1.font.color.rgb = RGBColor(0, 0, 0)

        if slide.slide_layout.name == "Objective":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Background Header PH"):
                shape.text = "Background"
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Actions Header PH"):
                shape.text = "Actions to be taken"
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Background Subheader PH"):
                shape.text = "What is the project objective?"
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Actions Subheader PH"):
                shape.text = "What decisions will this inspire?"
            if Document(brief_doc_path).tables[0].cell(0, 0).text == "Current Business Challenge":
                brief_objective_table = Document(brief_doc_path).tables[0]
                if shape.placeholder_format.idx == get_ph_idx(slide_type, "Background Text PH"):
                    shape.text = brief_objective_table.cell(1, 0).text
                if shape.placeholder_format.idx == get_ph_idx(slide_type, "Actions Text PH"):
                    shape.text = brief_objective_table.cell(1, 1).text

        if slide.slide_layout.name == "Scope":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Overview PH"):
                shape.text = "Overview"

        if slide.slide_layout.name == "Divider_normal":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Divider Text PH"):
                shape.text = "What is the next section covering?"

        if slide.slide_layout.name == "Scope_of_Work":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Scope Summary PH"):
                sow_cat = slide_df['Category'].unique()[0]
                sow_subcat = slide_df['Subcategory'].unique()[0]
                sow_country = slide_df['Country'].unique()[0]
                sow_segment = slide_df['Segment'].unique()[0]
                sow_period = slide_df['Cycle'].unique()[0]
                if sow_cat != sow_subcat:
                    sow_cat_subcat = f"{sow_cat} ({sow_subcat})"
                else:
                    sow_cat_subcat = sow_cat
                shape.text = f"We searched for a unique sample set talking about {sow_cat_subcat} by {sow_segment} in {sow_country} during {sow_period}."

        if (slide.slide_layout.name == "Content_Vertical") | (slide.slide_layout.name == "Content_Horizontal_Chart"):
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Country PH"):
                shape.text = slide_df['Country'].astype(str).unique()[0]
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Segment PH"):
                ## Single Segment
                segment_simple = None
                if "Segment" in slide_df.columns:
                    segment_simple = str([x.split(" ")[0] for x in slide_df['Segment'].unique()][0])
                ## Cross Segment
                elif "Segment 1" in slide_df.columns:  ## to adjust
                    segment_simple = str([x.split(" ")[0] for x in slide_df['Segment 1'].unique()][0])
                shape.text = segment_simple
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Core Commentary PH"):
                shape.text = "For core commentary."
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Supplementary Commentary PH"):
                shape.text = "For supplementary commentary."
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Footnote PH"):
                if slide_type == "Drivers":
                    num_ppl_all = int(slide_df['No_of_People'].unique()[0])
                    shape.text = f"Scores are derived from conversations by stated sample set (n={num_ppl_all:,}) and sorted in descending order."
                elif slide_type in ["Performance", "Emerging Trends"]:
                    shape.text = "Scores are derived from conversations by stated sample set (n=) and sorted in descending order."
                else:
                    pass
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Questions Text PH"):
                try:
                    questions_to_answer = str(slide_df['Questions'].unique()[0])
                    shape.text = questions_to_answer
                except:
                    shape.text = "<Pls copy 'Questions to answer' from AI Brief>"

        if slide.slide_layout.name == "KOL":
            if shape.placeholder_format.idx == get_ph_idx(slide_type, "Title PH"):
                shape.text = "Who are the top influencers/Key Opinion Leaders?"


def _set_cell_border(cell, ln_type, border_color="0000FF", border_width="101600"):
    """ Hack function to enable the setting of border width and border color
        - bottom border only at present
        (c) Steve Canny
    """
    def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    if ln_type == "none":
        return cell
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

##    lnR = SubElement(
##        tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
##    solidFill = SubElement(lnR, 'a:solidFill')
##    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)

    if ln_type == "right":
        # Right Cell Border
        lnR = SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        lnR_solidFill = SubElement(lnR, 'a:solidFill')
        lnR_srgbClr = SubElement(lnR_solidFill, 'a:srgbClr', val=border_color)
        lnR_prstDash = SubElement(lnR, 'a:prstDash', val='solid')
        lnR_round_ = SubElement(lnR, 'a:round')
        lnR_headEnd = SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
        lnR_tailEnd = SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')
        ln_type_config = lnR
        solidFill = SubElement(ln_type_config, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    elif ln_type == "left":
        # Left Cell Border
        lnL = SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        lnL_solidFill = SubElement(lnL, 'a:solidFill')
        lnL_srgbClr = SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
        lnL_prstDash = SubElement(lnL, 'a:prstDash', val='solid')
        lnL_round_ = SubElement(lnL, 'a:round')
        lnL_headEnd = SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
        lnL_tailEnd = SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')
        ln_type_config = lnL
        solidFill = SubElement(ln_type_config, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    elif ln_type == "top":
        # Top Cell Border
        lnT = SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        lnT_solidFill = SubElement(lnT, 'a:solidFill')
        lnT_srgbClr = SubElement(lnT_solidFill, 'a:srgbClr', val=border_color)
        lnT_prstDash = SubElement(lnT, 'a:prstDash', val='solid')
        lnT_round_ = SubElement(lnT, 'a:round')
        lnT_headEnd = SubElement(lnT, 'a:headEnd', type='none', w='med', len='med')
        lnT_tailEnd = SubElement(lnT, 'a:tailEnd', type='none', w='med', len='med')
        ln_type_config = lnT
        solidFill = SubElement(ln_type_config, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    elif ln_type == "bottom":
        # Bottom Cell Border
        lnB = SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        lnB_solidFill = SubElement(lnB, 'a:solidFill')
        lnB_srgbClr = SubElement(lnB_solidFill, 'a:srgbClr', val=border_color)
        lnB_prstDash = SubElement(lnB, 'a:prstDash', val='solid')
        lnB_round_ = SubElement(lnB, 'a:round')
        lnB_headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
        lnB_tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')
        ln_type_config = lnB
        solidFill = SubElement(ln_type_config, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    elif ln_type == "all":
        # All Cell Borders
        # recursion
        _set_cell_border(cell=cell, ln_type="left", border_color=border_color, border_width=border_width)
        _set_cell_border(cell=cell, ln_type="right", border_color=border_color, border_width=border_width)
        _set_cell_border(cell=cell, ln_type="top", border_color=border_color, border_width=border_width)
        _set_cell_border(cell=cell, ln_type="bottom", border_color=border_color, border_width=border_width)

    return cell


def format_table_cells(table, table_format):
    """
        Apply table/cell formatting as required //
        WORK IN PROGRESS: ALLOW FOR ROW-WISE/COL-WISE SPECIAL FORMATTING
    """
    ## General formatting - iterate over all cells

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            ## Set cell border (XML)
            cell = _set_cell_border(cell, table_format['cell_border'], table_format['cell_border_colour'],
                                    table_format['cell_border_size'])
            cell.vertical_anchor = table_format['cell_vert_anchor']
            ## Set fill type to transparent or solid color
            if table_format['cell_fill_transparent']:
                cell.fill.background()
            else:
                ## Set foreground (fill) color to a specific RGB color
                cell.fill.solid()
                cell.fill.fore_color.rgb = table_format['cell_fill_colour']
            ## Set text alignment left/center/right
            for para in cell.text_frame.paragraphs:
                para.alignment = table_format['cell_para_align']
                ## Set font Type / Size / Colour / Bold / Italics
                for run in para.runs:
                    run.font.name = table_format['cell_font_name']
                    run.font.size = table_format['cell_font_size']
                    run.font.color.rgb = table_format['cell_font_colour']
                    run.font.bold = table_format['cell_font_bold']
                    run.font.italic = table_format['cell_font_italic']

    def cell_para_run_font(cell):
        if not table_format['cell_fill_transparent']:
            cell.fill.fore_color.rgb = cell_format['cell_fill_colour']
        for para in cell.text_frame.paragraphs:
            para.alignment = cell_format['cell_para_align']
            for run in para.runs:
                run.font.size = cell_format['cell_font_size']
                run.font.color.rgb = cell_format['cell_font_colour']
                run.font.bold = cell_format['cell_font_bold']
                run.font.italic = cell_format['cell_font_italic']

    ## Slide-specific formatting - iterate over specified cells in 'special_format'
    if len(table_format['special_format']) > 0:
        for cell_format in table_format['special_format']:
            (cell_row, cell_col) = cell_format['cell_position']
            if type(cell_row) == type(cell_col) == int:
                cell = table.cell(cell_row, cell_col)
                table.rows[cell_row].height = cell_format['cell_height']
                table.columns[cell_col].width = cell_format['cell_width']
                cell_para_run_font(cell)

            if cell_col == "row":
                for row_idx, row in enumerate(table.rows):
                    if row_idx == cell_row:
                        for cell_idx, cell in enumerate(row.cells):
                            table.columns[cell_idx].width = cell_format['cell_width']
                            row.height = cell_format['cell_height']
                            cell_para_run_font(cell)
            if cell_row == "col":
                for row in table.rows:
                    for cell_idx, cell in enumerate(row.cells):
                        if cell_idx == cell_col:
                            table.columns[cell_idx].width = cell_format['cell_width']
                            row.height = cell_format['cell_height']
                            cell_para_run_font(cell)

    else:
        pass


def get_table_format(slide_df, table_config_json, brief_doc_path):
    """
        Return table/cell formatting as specified //
        WORK IN PROGRESS: UPDATED DRIVERS AND CROSS SEGMENT
    """
    # slide_type = slide_df['Slide_Type'].unique()[0]
    slide_type="Performance"
    table_format_dict = {
        # "Project_Scope": {
        #     "row_end_index": len(Document(brief_doc_path).tables[1].rows),
        #     "col_end_index": 2,
        #     "cell_fill_colour": RGBColor(242, 242, 242),
        #     "cell_font_colour": RGBColor(0, 0, 0),
        #     "cell_font_size": Pt(14),
        #     "cell_font_name": 'Poppins',
        #     "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
        #     "cell_vert_anchor": MSO_VERTICAL_ANCHOR.MIDDLE,
        #     "cell_border": "none",
        #     "cell_border_colour": "FFFFFF",
        #     "cell_border_size": "50800",
        #     "cell_font_bold": True,
        #     "cell_font_italic": False,
        #     "cell_fill_transparent": False,
        #     "special_format": [
        #         {
        #             "cell_position": ("col", 0),
        #             "cell_fill_colour": RGBColor(242, 242, 242),
        #             "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
        #             "cell_font_colour": RGBColor(0, 0, 0),
        #             "cell_font_size": Pt(14),
        #             "cell_font_bold": True,
        #             "cell_font_italic": False,
        #             "cell_height": Inches(0.4),
        #             "cell_width": Inches(2.5)
        #         },
        #         {
        #             "cell_position": ("col", 1),
        #             "cell_fill_colour": RGBColor(255, 255, 255),
        #             "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
        #             "cell_font_colour": RGBColor(0, 0, 0),
        #             "cell_font_size": Pt(14),
        #             "cell_font_bold": False,
        #             "cell_font_italic": False,
        #             "cell_height": Inches(0.4),
        #             "cell_width": Inches(10)
        #         },
        #     ]
        # },
        # "SOW": {
        #     "row_end_index": 5,
        #     "col_end_index": 1,
        #     "cell_fill_colour": RGBColor(242, 242, 242),
        #     "cell_font_colour": RGBColor(0, 0, 0),
        #     "cell_font_size": Pt(11),
        #     "cell_font_name": 'Poppins',
        #     "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
        #     "cell_vert_anchor": MSO_VERTICAL_ANCHOR.MIDDLE,
        #     "cell_border": "none",
        #     "cell_border_colour": "FFFFFF",
        #     "cell_border_size": "50800",
        #     "cell_font_bold": False,
        #     "cell_font_italic": False,
        #     "cell_fill_transparent": True,
        #     "special_format": [
        #         {
        #             "cell_position": (0, 0),
        #             "cell_fill_colour": RGBColor(255, 255, 255),
        #             "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
        #             "cell_font_colour": RGBColor(0, 0, 0),
        #             "cell_font_size": Pt(25),
        #             "cell_font_bold": True,
        #             "cell_font_italic": False,
        #             "cell_height": Inches(0.1),
        #             "cell_width": Inches(4)
        #         },
        #         {
        #             "cell_position": (1, 0),
        #             "cell_fill_colour": RGBColor(255, 255, 255),
        #             "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
        #             "cell_font_colour": RGBColor(0, 0, 0),
        #             "cell_font_size": Pt(60),
        #             "cell_font_bold": True,
        #             "cell_font_italic": False,
        #             "cell_height": Inches(0.1),
        #             "cell_width": Inches(4)
        #         },
        #         {
        #             "cell_position": (3, 0),
        #             "cell_fill_colour": RGBColor(255, 255, 255),
        #             "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
        #             "cell_font_colour": RGBColor(0, 0, 0),
        #             "cell_font_size": Pt(60),
        #             "cell_font_bold": True,
        #             "cell_font_italic": False,
        #             "cell_height": Inches(0.1),
        #             "cell_width": Inches(4)
        #         },
        #     ]
        # },
        "Performance": {
            "row_end_index": len(slide_df) + 2,
            "col_end_index": len([col for col in slide_df.columns if re.match('Segment *', col)]) + 1,
            "cell_fill_colour": RGBColor(242, 242, 242),
            "cell_font_colour": RGBColor(0, 0, 0),
            "cell_font_size": Pt(16),
            "cell_font_name": 'Poppins',
            "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
            "cell_vert_anchor": MSO_VERTICAL_ANCHOR.MIDDLE,
            "cell_border": "all",
            "cell_border_colour": "FFFFFF",
            "cell_border_size": "50800",
            "cell_font_bold": False,
            "cell_font_italic": False,
            "cell_fill_transparent": False,
            "special_format": [
                {
                    "cell_position": (0, "row"),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.5),
                    "cell_width": Inches(2)
                },
                {
                    "cell_position": (1, "row"),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(2)
                },
                {
                    "cell_position": ("col", 0),
                    "cell_fill_colour": RGBColor(242, 242, 242),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(16),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(5)
                },
                {
                    "cell_position": (0, 0),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(18),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(5)
                },
                {
                    "cell_position": (1, 0),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(5)
                },
            ]
        },
        "Drivers": {
            "row_end_index": len(slide_df) + 1,
            "col_end_index": len(slide_df.columns) - 11,
            "cell_fill_colour": RGBColor(255, 255, 255),
            "cell_font_colour": RGBColor(0, 0, 0),
            "cell_font_size": Pt(10),
            "cell_font_name": 'Poppins',
            "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
            "cell_vert_anchor": MSO_VERTICAL_ANCHOR.MIDDLE,
            "cell_border": "all",
            "cell_border_colour": "FFFFFF",
            "cell_border_size": "50800",
            "cell_font_bold": False,
            "cell_font_italic": False,
            "cell_fill_transparent": False,
            "special_format": [
                {
                    "cell_position": (0, "row"),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": True,
                    "cell_font_italic": True,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(0.5)
                },
                {
                    "cell_position": ("col", 0),
                    "cell_fill_colour": RGBColor(217, 217, 217),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(0.75)
                },
                {
                    "cell_position": ("col", 1),
                    "cell_fill_colour": RGBColor(242, 242, 242),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(4)
                },
                {
                    "cell_position": (0, 0),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(18),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(1)
                },
                {
                    "cell_position": (0, 1),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(18),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(4)
                },
            ]
        },
        "KOL": {
            "row_end_index": len(slide_df.columns) - 1,
            "col_end_index": 6,
            "cell_fill_colour": RGBColor(255, 255, 255),
            "cell_font_colour": RGBColor(0, 0, 0),
            "cell_font_size": Pt(10),
            "cell_font_name": 'Poppins',
            "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
            "cell_vert_anchor": MSO_VERTICAL_ANCHOR.MIDDLE,
            "cell_border": "none",
            "cell_border_colour": "FFFFFF",
            "cell_border_size": "50800",
            "cell_font_bold": False,
            "cell_font_italic": False,
            "cell_fill_transparent": True,
            "special_format": [
                {
                    "cell_position": ("col", 0),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(0, 0, 0),
                    "cell_font_size": Pt(12),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(2.1)
                },
                {
                    "cell_position": (0, "row"),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.CENTER,
                    "cell_font_colour": RGBColor(127, 127, 127),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": False,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(2.1)
                },
                {
                    "cell_position": (0, 0),
                    "cell_fill_colour": RGBColor(255, 255, 255),
                    "cell_para_align": PP_PARAGRAPH_ALIGNMENT.LEFT,
                    "cell_font_colour": RGBColor(127, 127, 127),
                    "cell_font_size": Pt(10),
                    "cell_font_bold": True,
                    "cell_font_italic": False,
                    "cell_height": Inches(0.1),
                    "cell_width": Inches(2.1)
                },
            ]
        }
    }
    with open(table_config_json, 'w') as json_file:
        json.dump(table_format_dict, json_file)

    return table_format_dict[slide_type]


def fill_table_data(slide, slide_df, slide_master_df,base=None):
    """
        Fill Table Placeholders with required inputs //
        Apply table formatting //
        WORK IN PROGRESS: UPDATE FOR CROSS SEGMENT, UPDATE KOL, EXTRACT PROJECT SCOPE FROM AI BRIEF
    """
    # slide_type = slide_df['Slide_Type'].unique()[0]
    slide_type="Performance"
    for shape in slide.shapes:
        if (shape.is_placeholder) & ("Table" in shape.name):
            new_table_format = get_table_format(slide_df, table_config_json, brief_doc_path)
            new_shape = shape.insert_table(new_table_format['row_end_index'], new_table_format['col_end_index'])
            tbl = new_shape._element.graphic.graphicData.tbl
            style_id = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
            tbl[0][-1].text = style_id
            new_table = new_shape.table

            if (slide.slide_layout.name == "Content_Vertical") & (slide_type == "Performance"):
                ## Single Segment
                # if "Segment" in slide_df.columns:
                #     header_top_left_cell = new_table.cell(0, 0)
                #     header_top_left_cell.text = slide_df['Type'].astype(str).unique()[0] + "/" + \
                #                                 slide_df['Subtype'].astype(str).unique()[0]
                #     header_top_right_cell = new_table.cell(0, 1)
                #     header_top_right_cell.text = slide_df['Segment'].astype(str).unique()[0]
                #     header_bot_right_cell = new_table.cell(1, 1)
                #     segment_num_ppl = int(slide_df['No_of_People'].unique()[0])
                #     header_bot_right_cell.text = f"(n={segment_num_ppl:,})"
                #     for row_index, row in enumerate(new_table.rows):
                #         if row_index > 1:
                #             for cell_index, cell in enumerate(row.cells):
                #                 if cell_index == 0:
                #                     cell.text = slide_df['Content'].astype(str).iloc[row_index - 2]
                #                 if cell_index == 1:
                #                     cell.text = slide_df['Measure Value'].astype(str).iloc[row_index - 2]
                ## Cross Segment
                # if "Segment 1" in slide_df.columns:  ## to adjust
                header_top_left_cell = new_table.cell(0, 0)
                header_top_left_cell.text = "Type" + "/" + "Subtype"
                num_segments = len([col for col in slide_df.columns if re.match('Segment *', col)])
                for r in range(1, num_segments + 1):
                    header_segment_cell = new_table.cell(0, r)
                    if r ==1:
                        header_segment_cell.text = f'Segment/Brand {r} \n (Benchmark)'  ## to adjust
                    else:
                        header_segment_cell.text = f'Segment/Brand {r}'  ## to adjust

                    subheader_num_cell = new_table.cell(1, r)
                    # segment_num_ppl = int(slide_df[f'No_of_People Segment {r}'].dropna().unique()[0])  ## to adjust
                    if base==None:
                        subheader_num_cell.text = f"(n= )"
                    else:
                        subheader_num_cell.text = f"(n= {str(base)} )"

                    for row_index, row in enumerate(new_table.rows):
                        if row_index > 1:
                            for cell_index, cell in enumerate(row.cells):
                                if cell_index == 0:
                                    try:
                                        cell.text = slide_df['Content'].astype(str).iloc[row_index - 2]
                                    except:
                                        pass
                                if cell_index == r:
                                    
                                    cell_value = slide_df[f'Segment/Brand {r}'].astype(str).iloc[row_index - 2]
                                  
                                    if cell_value == 'nan':
                                        cell.text = "-"
                                    else:
                                        cell.text = cell_value  ## to adjust


        ## Formats table for all slide types above except KOL slide which has 2 tables so nested below
            format_table_cells(new_table, new_table_format)


def fill_chart_data(slide, slide_df):
    """
        Fill Chart Placeholders with required inputs //
        WORK IN PROGRESS: PENDING EMERGING TRENDS LOGIC CONFIRMATION
    """
    temp_slide_type = slide_df['Slide_Type'].unique()[0]
    for shape in slide.shapes:
        if (shape.is_placeholder) & ("Chart" in shape.name):
            new_table_format = get_table_format(slide_df, brief_doc_path)
            slide_df = slide_df[['']]
            new_shape = shape.insert_chart(XL_CHART_TYPE.LINE_MARKERS, )
            new_table = new_shape.table
    pass


"""""""""""""""""""""""""""""""""""""""  Stat-Sig """""""""""""""""""""""""""""""""""""""""""""

def get_slide_table(slide):
    """
        Get table from slide already generated by earlier insert_table
    """
    for shape in slide.shapes:
        if "Table Placeholder" in shape.name:
            table = shape.table
            return table


def get_slide_table_values(table):
    """
        Get table data from slide already filled by earlier fill_table_data
    """
    table_df = pd.DataFrame()
    for row_idx, row in enumerate(table.rows):
        row_lst = []
        for cell_idx, cell in enumerate(row.cells):
            row_lst.append(cell.text)
        table_df = pd.concat((table_df, pd.DataFrame(row_lst).T),axis=0)
    table_df = table_df.reset_index(drop=True)
    return table_df


def find_threshold(number,base=None):
    """
        Get threshold value for segment / brand statistical significance comparison
    """
    if (base==None) or (base>1100):
        if number>45:
            return (4.1*1.25)
        elif number >30:
            return (1.76*1.25)
        elif number >15:
            return (0.98*1.25)
        elif number > 0:
            return (0.56*1.25)
        else:
            return "Error"
    elif (base>750):
        if number>45:
            return (4.7*1.25)
        elif number >30:
            return (2.21*1.25)
        elif number >15:
            return (1.73*1.25)
        elif number > 0:
            return (0.97*1.25)
        else:
            return "Error"
    elif (base>300):
        if number>45:
            return (6.3*1.25)
        elif number >30:
            return (3.22*1.25)
        elif number >15:
            return (2.2*1.25)
        elif number > 0:
            return (1.23*1.25)
        else:
            return "Error"
    else:
        return "Base less than 300!!"


def find_largest_and_second_largest(numbers_list):
    """
        Get first and second largest based on values of each row for Max Stat-sig analysis
    """
    if len(set(numbers_list)) <= 1:
        return None, None
    largest = max(numbers_list)
    numbers_list = [num for num in numbers_list if int(round(num,0))!= int(round(largest,0))]
    second_largest = max(numbers_list)
    return largest, second_largest


def get_benchmark_col_idx(slide_table, benchmark):
    """
        Get benchmark column name (int value) for Benchmark Stat-sig analysis
    """
    for cell_idx, cell in enumerate(slide_table.iloc[0]):
        if cell == benchmark:
            benchmark_col = cell_idx
            return benchmark_col
    print(f"Benchmark does not exactly match any of the segments/brands specified.")
    return None


def set_cell_statsig_colour(cell, colour):
    """
        Set cell font colour to Green/Red if Stat Significant/Superior/Inferior
    """
    RED = RGBColor(255,0,0)
    GREEN = RGBColor(0, 176, 80)
    try:
        # cell.text_frame.paragraphs[0].text= str(int(round(float( cell.text_frame.text),0)))+" %"
        ## 1.d.p instead of 0.d.p
        cell.text_frame.paragraphs[0].text= str(round(float( cell.text_frame.text),1))+" %"
    except:
        pass
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            if colour == "Green":
                run.font.color.rgb = GREEN
            if colour == "Red":
                run.font.color.rgb = RED


def apply_stat_sig(slide, slide_df,statsig_type,base):
    """
        Max Stat-sig: Compare data points of cross segments to get diff between 1st/2nd largest and determine stat-significance of largest value. //
        Benchmark Stat-sig: Compare each data point to ref benchmark value to determine stat-superior/inferior of benchmark value against other value. //
        Set cell font colour to green/red if diff exceeds Stat-Sig threshold.
    """


    ## Get base number of people
 
    # performance_type = slide_df['Type'].unique()[0]
    # performance_subtype = slide_df['Subtype'].unique()[0]

    ## Get table from slide
    table = get_slide_table(slide)
    slide_table = get_slide_table_values(table)


    # benchmark = "Segment/Brand 1"
    if statsig_type == "max":
        ## Get 1st/2nd largest, threshold
        for table_idx, row in slide_table.iterrows():
            if table_idx > 1:
                num_list = list(row[1:len(row)])
                # num_list = [float(num.strip('%')) for num in num_list if num != "-"]
                num_list = [0 if num =="-" else float(num) for num in num_list]

                (largest, second) = find_largest_and_second_largest(num_list)
                if (largest != None) & (second != None):
                    largest_index_list = [num_index for num_index, num in enumerate(num_list) if int(round(num,0)) == int(round(largest,0))]
                # second_index = num_list.index(second)
                threshold = find_threshold(largest, base)
                diff = largest - second
                print(largest_index_list)
                print(diff)
                print(threshold)
                for num_idx, num in enumerate(num_list):
                    cell = table.cell(table_idx, num_idx + 1)

                    if ((num_idx in largest_index_list) & (diff>threshold)):

                        set_cell_statsig_colour(cell, "Green")
                    else:
                        set_cell_statsig_colour(cell, None)

                    # for largest_index in largest_index_list:
                    #     second_index = num_list.index(second)
                    #     threshold = find_threshold(largest, base)
                    #     diff = largest - second

                    #     cell = table.cell(table_idx, largest_index + 1)
                    #     if diff > threshold:
                    #         set_cell_statsig_colour(cell, "Green")
                    #     else:
                    #         set_cell_statsig_colour(cell, None)

    if statsig_type=="benchmark" :
        ## Get benchmark segment / respective column
        # ref_benchmark = benchmark

        ## set row index and column index to start iterating from
        # if "segment" in statsig_type:
        row_start_index = 2
        col_start_index = 1
        # if "brand" in statsig_type:
        #     row_start_index = 1
        #     col_start_index = 2

        ## Get threshold and compare diff of each segment value against benchmark segment value
        for table_idx, row in slide_table.iterrows():
            if table_idx >= row_start_index:
                num_list = list(row[col_start_index:len(row)])
                # num_list = [float(num.strip("%")) for num in num_list if num != "-"]
                num_list = [0 if num =="-" else float(num) for num in num_list]

                # benchmark_idx = get_benchmark_col_idx(slide_table, ref_benchmark) - col_start_index
                benchmark_idx = 0

                if benchmark_idx < len(num_list):
                    benchmark_value = num_list[benchmark_idx]
                    base_num=base
                    # if "segment" in statsig_type:
                    #     base_num = base_list[benchmark_idx]
                    # if "brand" in statsig_type:
                    #     base_num = slide_df['No_of_People'].dropna().unique()[0]
                    threshold = find_threshold(benchmark_value, base_num)
                    # table.cell(table_idx, benchmark_idx + col_start_index).text_frame.paragraphs[0].runs[
                    #     0].font.bold = True
                    for num_idx, num in enumerate(num_list):
                        
                        diff = num - benchmark_value
                        cell = table.cell(table_idx, num_idx + col_start_index)

                        if num ==0:
                            set_cell_statsig_colour(cell, None)

                        else:
                        ## default logic, i.e. other value is stat-larger than benchmark value shall be red
                            try:
                                if (diff > threshold) & ("inverse" not in statsig_type):
                                    set_cell_statsig_colour(cell, "Red")
                                if (diff < - threshold) & ("inverse" not in statsig_type):
                                    set_cell_statsig_colour(cell, "Green")
                                ## inverse logic, i.e. other value is stat-larger than benchmark value shall be green
                                if (diff > threshold) & ("inverse" in statsig_type):
                                    set_cell_statsig_colour(cell, "Green")
                                if (diff < - threshold) & ("inverse" in statsig_type):
                                    set_cell_statsig_colour(cell, "Red")
                                else :
                                    set_cell_statsig_colour(cell, None)
                            except:
                                set_cell_statsig_colour(cell, None)

                    table.cell(table_idx, benchmark_idx + col_start_index).text_frame.paragraphs[0].runs[
                        0].font.bold = True

    """ Try to change to 1 d.p,fail"""
    # for row_idx, row in enumerate(table.rows):
    #     if row_idx > 1:
    #         for cell_idx,cell in enumerate(row.cells):
    #             if cell_idx>0:
    #                 if cell.text_frame.paragraphs[0].text != "-":
    #                     cell.text_frame.paragraphs[0].text= str(int(round(float( cell.text_frame.text),0)))
          
                        # Access the text in the cell


        # num_list = list(row[col_start_index:len(row)])
        # num_list = [float(num.strip("%")) for num in num_list if num != "-"]
        # for num_idx, num in enumerate(num_list):

  
        # for para in cell.text_frame:
        #     para.text= str(int(round(float(para.text),0)))
        # cell.text_frame.text= str(int(round(float( cell.text_frame.text),0)))


      
    return slide_table


"""""""""""""""""""""""""""""""""""""""  Execute Functions """""""""""""""""""""""""""""""""""""""""""""


def main_execute(df,statsig_type,base,ref_benchmark):
    """
        Execute all functions to generate sequential data dict and PPT slide sequence with text/table inputs //
        Print for each slide added //
        Save output file with Arg
    """
    prs = Presentation(default_template)

    # slide_sequence_dict = get_slide_sequence_values(mapping_path)
    slide_master_df = get_slide_master_df(default_template, placeholder_description_input,
                                          slide_master_placeholders_output)
    
    # print(slide_sequence_dict)

    temp_slide_layout_index = slide_master_df[slide_master_df['Slide_Type'] == "Performance"]['Slide_Layout_Index'].unique()[0]
    new_slide = prs.slides.add_slide(prs.slide_layouts[temp_slide_layout_index])
    # print(f"{slide_key} slide added using {new_slide.slide_layout.name}")
    # fill_text_data(new_slide, df, slide_master_df)
    fill_table_data(new_slide, df, slide_master_df,base=base)
    apply_stat_sig (new_slide, df,statsig_type,base)
    # if 'Statsig_Type' in slide_df.columns:
    #     apply_stat_sig(new_slide, slide_df)

    prs.save(output_file_name)
    ppt_buffer=io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer
# main_execute(output_file_name)